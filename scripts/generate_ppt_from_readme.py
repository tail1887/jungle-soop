import re
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
README = ROOT / "README.md"
ASSETS_DIR = ROOT / "artifacts" / "ppt_assets"
OUTPUT_PPT = ROOT / "artifacts" / "jungle-soop-presentation.pptx"


def extract_mermaid_blocks(readme_text: str) -> list[str]:
    pattern = r"```mermaid\s*\n(.*?)```"
    return re.findall(pattern, readme_text, flags=re.DOTALL)


def render_mermaid(block: str, name: str) -> Path:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    mmd_path = ASSETS_DIR / f"{name}.mmd"
    png_path = ASSETS_DIR / f"{name}.png"
    mmd_path.write_text(block, encoding="utf-8")

    cmd = (
        f'npx -y @mermaid-js/mermaid-cli '
        f'-i "{mmd_path}" -o "{png_path}" '
        f'-t default -b white --scale 2'
    )
    subprocess.run(cmd, check=True, cwd=ROOT, shell=True)
    return png_path


def add_title(slide, title: str, subtitle: str = ""):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(1.2))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(32, 43, 76)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(12), Inches(0.8))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(18)
        sp.font.color.rgb = RGBColor(75, 85, 99)


def add_bullets(slide, bullets: list[str], left=0.9, top=2.0, width=5.8, height=4.5):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(31, 41, 55)
        p.space_after = Pt(12)


def add_image(slide, image_path: Path, left=6.7, top=1.8, width=6.1, height=4.9):
    slide.shapes.add_picture(
        str(image_path),
        Inches(left),
        Inches(top),
        width=Inches(width),
        height=Inches(height),
    )


def add_center_message(slide, text: str):
    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1.8))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(17, 24, 39)


def build_ppt(diagrams: dict[str, Path]):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Cover
    s = prs.slides.add_slide(blank)
    add_title(s, "Jungle-Soop", "기숙사 소규모 모임 커뮤니티 서비스 발표")
    add_center_message(s, "빠르게 만들고, 가볍게 참여하는 생활형 모임 플랫폼")

    # 2. Problem
    s = prs.slides.add_slide(blank)
    add_title(s, "문제 정의", "왜 이 서비스를 만들었나?")
    add_bullets(
        s,
        [
            "• 모임 정보가 단톡방/지인 네트워크에 분산됨",
            "• 신규 사용자는 참여 기회를 찾기 어려움",
            "• 소규모 모임 생성/참여 절차가 번거로움",
        ],
    )
    add_image(s, diagrams["frontend_flow"], left=6.6, top=1.8, width=6.2, height=4.9)

    # 3. Goal
    s = prs.slides.add_slide(blank)
    add_title(s, "서비스 목표", "핵심 가치 3가지")
    add_bullets(
        s,
        [
            "• 발견성: 열려 있는 모임을 빠르게 탐색",
            "• 즉시성: 최소 입력으로 모임 개설",
            "• 지속성: 참여 이후 상호작용 경험 강화",
        ],
    )
    add_image(s, diagrams["frontend_login_wire"], left=6.6, top=2.0, width=6.0, height=4.4)

    # 4. User Flow
    s = prs.slides.add_slide(blank)
    add_title(s, "핵심 사용자 흐름", "로그인부터 참여까지")
    add_bullets(
        s,
        [
            "• 로그인",
            "• 모임 목록 탐색",
            "• 모임 상세 확인",
            "• 참여 신청",
            "• 참여 상태 확인",
        ],
    )
    add_image(s, diagrams["frontend_flow"], left=6.2, top=1.8, width=6.6, height=4.9)

    # 5. Feature Snapshot
    s = prs.slides.add_slide(blank)
    add_title(s, "주요 기능", "현재 구현 범위")
    add_bullets(
        s,
        [
            "• 회원가입/로그인 (JWT 기반 인증)",
            "• 모임 생성/수정/삭제",
            "• 모임 목록/상세 조회",
            "• 모임 참여/취소 및 정원 제어",
            "• 반응형 UI + AJAX 상호작용",
        ],
    )
    add_image(s, diagrams["frontend_list_wire"], left=6.6, top=1.9, width=6.1, height=4.8)

    # 6. Overall Architecture
    s = prs.slides.add_slide(blank)
    add_title(s, "전체 아키텍처", "Client - Application - Data - Infra")
    add_bullets(
        s,
        [
            "• Client: HTML/CSS/Tailwind + JavaScript/jQuery/AJAX",
            "• Application: Flask + Jinja2 + REST API",
            "• Data: MongoDB",
            "• Infra: Docker on AWS EC2 + GitHub Actions",
        ],
    )
    add_image(s, diagrams["overall_arch"], left=6.0, top=1.5, width=6.9, height=5.5)

    # 7. Frontend Architecture
    s = prs.slides.add_slide(blank)
    add_title(s, "프론트엔드 아키텍처", "SSR + 페이지 스크립트 + AJAX")
    add_bullets(
        s,
        [
            "• SSR 응답으로 초기 화면 로딩",
            "• 페이지 JS 모듈에서 이벤트 처리",
            "• AJAX로 데이터 요청 후 DOM 갱신",
        ],
    )
    add_image(s, diagrams["frontend_linear"], left=6.3, top=2.0, width=6.4, height=3.7)

    # 8. Backend Architecture
    s = prs.slides.add_slide(blank)
    add_title(s, "백엔드 아키텍처", "Route -> Service -> Repository")
    add_bullets(
        s,
        [
            "• Route: HTTP/인증/응답 책임",
            "• Service: 도메인 규칙 및 유스케이스",
            "• Repository: MongoDB I/O 분리",
            "• 계층 분리로 테스트 및 유지보수성 강화",
        ],
    )
    add_image(s, diagrams["backend_linear"], left=6.1, top=1.9, width=6.7, height=4.4)

    # 9. Technical Challenge
    s = prs.slides.add_slide(blank)
    add_title(s, "기술 설명", "해결한 핵심 문제")
    add_bullets(
        s,
        [
            "• 인증/권한 제어 흐름 표준화",
            "• 참여 로직에서 중복/정원 초과 충돌 처리",
            "• SSR + API 응답 모델 재사용으로 일관성 확보",
        ],
    )
    add_image(s, diagrams["seq_login"], left=6.1, top=1.8, width=6.7, height=4.8)

    # 10. Roadmap
    s = prs.slides.add_slide(blank)
    add_title(s, "확장 계획", "다음 단계")
    add_bullets(
        s,
        [
            "• 단기: 댓글/참여 알림 기능",
            "• 중기: 관심사 기반 모임 추천",
            "• 장기: 일정 연동 및 사용자 평판 시스템",
            "• 목표: 재방문율과 참여 지속성 향상",
        ],
    )
    add_image(s, diagrams["seq_join"], left=6.2, top=1.8, width=6.6, height=4.8)

    OUTPUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPT))


def main():
    readme_text = README.read_text(encoding="utf-8")
    blocks = extract_mermaid_blocks(readme_text)

    index_map = {
        "overall_arch": 0,
        "frontend_linear": 1,
        "backend_linear": 2,
        "frontend_flow": 3,
        "frontend_login_wire": 4,
        "frontend_list_wire": 5,
        "seq_login": 8,
        "seq_join": 11,
    }

    diagrams = {}
    for key, idx in index_map.items():
        diagrams[key] = render_mermaid(blocks[idx], key)

    build_ppt(diagrams)
    print(f"Generated: {OUTPUT_PPT}")


if __name__ == "__main__":
    main()
